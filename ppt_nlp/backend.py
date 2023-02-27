import os
import openai
import collections.abc
import pptx


def load_ppt(ppt):
    """
    Load a PowerPoint file and extract the text from each slide.

    Args:
    - ppt (str): The path to the PowerPoint file.

    Returns:
    - str: A string that contains the text from all the slides in the PowerPoint file.

    Raises:
    - FileNotFoundError: If the specified PowerPoint file is not found.
    - pptx.exc.PackageNotFoundError: If the specified file is not a PowerPoint file.

    """
    try:
        # Load the PowerPoint file
        pptx_file = pptx.Presentation(ppt)

        # Extract the text from each slide and concatenate it into a single string
        text = ""
        for slide in pptx_file.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text

        # Return the text
        return text

    except FileNotFoundError:
        print(f"Error: The specified file '{ppt}' was not found.")
        return None

    except pptx.exc.PackageNotFoundError:
        print(f"Error: The specified file '{ppt}' is not a valid PowerPoint file.")
        return None


def summarization(text):
    """
    Submit a question to the OpenAI API and receive a summarization of the answer.

    Args:
    - text (str): The question to submit to the OpenAI API.

    Returns:
    - str: A summarization of the answer to the question.

    Raises:
    - openai.Error: If there is an error communicating with the OpenAI API.

    """
    try:
        # Set the OpenAI API key and prompt
        openai.api_type = os.getenv("OPENAI_TYPE")
        openai.api_key = os.getenv("OPENAI_API_KEY")
        openai.api_base = os.getenv("OPENAI_BASE")
        openai.api_version = os.getenv("OPENAI_VERSION")
        prompt = text + "'tl;dr:"

        # Submit the question to the OpenAI API and receive a response
        response = openai.Completion.create(
            prompt=prompt,
            temperature=0,
            max_tokens=300,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0,
            engine="text-davinci-002",
        )

        # Extract the summarization from the response
        result = response["choices"][0]["text"].strip(" \n")

        # Return the summarization
        return result

    except BaseException as e:
        print(f"Error: {e}")
        return None
